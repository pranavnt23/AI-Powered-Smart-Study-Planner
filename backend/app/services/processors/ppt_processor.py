from pathlib import Path
from pptx import Presentation


class PPTProcessor:
    @staticmethod
    def extract_text(file_path: str) -> dict:
        try:
            presentation = Presentation(file_path)
            extracted_lines = []

            def append_text(text: str):
                if not text:
                    return
                cleaned = text.strip()
                if cleaned:
                    extracted_lines.append(cleaned)

            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                append_text(cell.text)

                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            append_text(paragraph.text)

                if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
                    append_text(slide.notes_slide.notes_text_frame.text)

            extracted_text = "\n".join(extracted_lines)

            return {
                "status": True,
                "file_type": "pptx",
                "file_name": Path(file_path).name,
                "content": extracted_text,
            }

        except Exception as error:
            return {
                "status": False,
                "file_type": "pptx",
                "file_name": Path(file_path).name,
                "content": "",
                "error": str(error),
            }
